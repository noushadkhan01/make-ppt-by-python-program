import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
#%matplotlib inline
from pptx import Presentation
from pptx.util import Inches
import os
def addBulletSlide(ppt_name):
    try:
        prs = Presentation(ppt_name + '.pptx')
    except:
        prs = Presentation()
    
    heading = input('Please Enter heading of Bullet Slide:-- ')
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = heading
    dct = dict(zip((1, 2, 3, 4, 5, 6, 7, 8, 9, 10), 
                   ('1st', '2nd', '3rd', '4th', '5th', '6th', '7th', '8th', '9th', '10th')))
    count = 1
    while True:
        first_text = input(f'Want to add {dct.get(count)} heading Please enter:--(if not type quit):-- ')
        if first_text.lower() != 'quit' and first_text.lower() != 'break':
            tf = body_shape.text_frame
            tf.text = first_text
            tf.lvel = 1
            count1 = 1
            while True:
                bullet_line = input(f'Want to add {dct.get(count1)} Subheading Please enter:--(if not type quit):-- ')
                if bullet_line.lower() != 'quit' and bullet_line.lower() != 'break':
                    p = tf.add_paragraph()
                    p.text = bullet_line
                    p.level = 1
                    count1 += 1
                    continue
                break
            count += 1
            continue
        break
    prs.save(ppt_name + '.pptx')
                

def addImageSlide(img_path, ppt_name):
    try:
        prs = Presentation(ppt_name + '.pptx')
    except:
        prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    left =Inches(1.6)
    top = Inches(1)
    pic = slide.shapes.add_picture(img_path, left, top)
    prs.save(ppt_name + '.pptx')
    os.remove(img_path)

def createppt(ppt_name, *args):
    #1 create text slide
    x = None
    d = input('Want to add bullet slide:--(y / n) ')
    if d.lower() == 'y' or d.lower() == 'yes':
        x = 1
    else:
        x = None
    if x:
        addBulletSlide(ppt_name)

    #2 create image slide
    for i in args:
        img_path = i + '.png'
        addImageSlide(img_path, ppt_name)
            
def makeppt(file_name, ppt_name):
    import matplotlib.pyplot as plt
    #%matplotlib inline
    try:
        with open(str(file_name)) as f:
            data = f.read()
            l = []
            for i in range(2):
                n1 = data.find('[')
                n2 = data.find(']')
                d = data[n1 + 1: n2]
                l.append([int(i) for i in d if i.isdigit()])
                data = data[n2+1:]
    except:
        return None
    x, y = l[0], l[1]
    plt.title('X vs Y dashed line plot ')
    plt.xlabel('x ---------------->')
    plt.ylabel('y ---------------->')
    f = plt.plot(x, y, color='green', marker='o', linestyle='dashed',
        linewidth=2, markersize=12)
    plt.grid()
    plt.savefig('plot.png', dpi = 100, pad_inches = 0.3)
    plt.close()
    plt.title('X vs Y Scatter plot')
    plt.xlabel('x ---------------->')
    plt.ylabel('y ---------------->')
    g = plt.scatter(x, y,color = ['r', 'g'], marker = 'o', edgecolor = 'orange', s = 150)
    plt.grid()
    plt.savefig('scatter.png', dpi = 100, pad_inches = 0.3)
    plt.close()
    return createppt(ppt_name, 'plot', 'scatter')

file_name = input('Enter the name of text file:-- ')
if '.txt' in file_name:
    file_name = file_name
else:
    file_name = file_name + '.txt'
name = input('enter name of pptx file:-- ')
makeppt(file_name, name)
